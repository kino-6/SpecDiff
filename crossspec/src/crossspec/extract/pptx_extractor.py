"""PPTX extractor."""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, List

from pptx import Presentation

from crossspec.claims import Authority
from crossspec.config import PptxConfig
from crossspec.extract.base import ExtractedClaim, Extractor


class PptxExtractor(Extractor):
    def __init__(self, path: Path, authority: Authority, config: PptxConfig) -> None:
        self.path = path
        self.authority = authority
        self.config = config

    def extract(self) -> Iterable[ExtractedClaim]:
        presentation = Presentation(self.path)
        for idx, slide in enumerate(presentation.slides, start=1):
            texts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text
                    if text:
                        texts.append(text)
            notes_text = ""
            if self.config.include_notes and slide.has_notes_slide:
                notes = slide.notes_slide
                if notes and notes.notes_text_frame:
                    notes_text = notes.notes_text_frame.text
            body = "\n".join(texts)
            if notes_text:
                body = f"[Slide {idx}]\n{body}\n(Notes)\n{notes_text}"
            else:
                body = f"[Slide {idx}]\n{body}"
            provenance = {"slide": idx}
            yield ExtractedClaim(
                text_raw=body,
                source_type="pptx",
                source_path=str(self.path),
                authority=self.authority,
                provenance=provenance,
            )
