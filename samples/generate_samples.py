"""Generate small sample artifacts for CrossSpec demo."""

from __future__ import annotations

import sys
from pathlib import Path
from typing import List

BASE_DIR = Path(__file__).resolve().parent
INPUT_DIR = BASE_DIR / "input"


def main() -> int:
    INPUT_DIR.mkdir(parents=True, exist_ok=True)
    _generate_eml()
    _generate_xlsx()
    _generate_pptx()
    _generate_pdf()
    return 0


def _generate_pdf() -> None:
    output_path = INPUT_DIR / "sample.pdf"
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
    except ModuleNotFoundError:
        print("reportlab not installed; skipping PDF generation.")
        return

    c = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter
    text = c.beginText(40, height - 50)
    paragraphs = [
        "Brake controller shall support safe deceleration under normal conditions.",
        "CAN interface must handle timing requirements for diagnostic frames.",
        "Error handling includes retry logic and fault logging for safety-critical paths.",
        "- Feature: brake control stability",
        "- Feature: can timing diagnostics",
        "- Feature: error handling coverage",
    ]
    for line in paragraphs:
        text.textLine(line)
    c.drawText(text)
    c.showPage()
    text = c.beginText(40, height - 50)
    text.textLine("Calibration settings must be retained in non-volatile memory.")
    text.textLine("Safety checks should run before enabling braking output.")
    c.drawText(text)
    c.save()
    print(f"Generated {output_path}")


def _generate_xlsx() -> None:
    output_path = INPUT_DIR / "sample.xlsx"
    try:
        from openpyxl import Workbook
    except ModuleNotFoundError:
        print("openpyxl not installed; skipping XLSX generation.")
        return

    wb = Workbook()
    ws = wb.active
    ws.title = "Q&A"
    ws.append(["Question", "Answer", "Status", "Owner"])
    rows = [
        ("How is brake torque limited?", "Via controller thresholds.", "Approved", "Alice"),
        ("Is CAN timing verified?", "Yes, per diagnostics.", "Approved", "Bob"),
        ("What is error handling retry count?", "3 attempts.", "Draft", "Chen"),
        ("Where is calibration stored?", "In NVM.", "Approved", "Dina"),
        ("Does safety check run at init?", "Yes, before output.", "Draft", "Eli"),
        ("Is diagnostics logging persistent?", "Stored in NVM.", "Approved", "Fay"),
        ("How is CAN bus load handled?", "Rate limiting.", "Draft", "Gus"),
        ("Brake response time?", "< 50ms.", "Approved", "Hana"),
        ("Error codes defined?", "See error handling section.", "Draft", "Ivy"),
        ("Any safety overrides?", "Yes, fail-safe mode.", "Approved", "Jun"),
    ]
    for row in rows:
        ws.append(row)
    wb.save(output_path)
    print(f"Generated {output_path}")


def _generate_pptx() -> None:
    output_path = INPUT_DIR / "sample.pptx"
    try:
        from pptx import Presentation
    except ModuleNotFoundError:
        print("python-pptx not installed; skipping PPTX generation.")
        return

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    def add_slide(title: str, bullets: List[str], notes: str | None = None) -> None:
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        for idx, bullet in enumerate(bullets):
            if idx == 0:
                body.text = bullet
            else:
                body.add_paragraph().text = bullet
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    add_slide(
        "Brake Feature Overview",
        ["Brake control stability", "Safety checks at init", "Diagnostics logging"],
        "Notes: emphasize braking safety and error handling.",
    )
    add_slide(
        "CAN Diagnostics",
        ["Timing requirements", "Error handling on bus", "Calibration updates"],
    )
    add_slide(
        "NVM and Calibration",
        ["Persistent storage", "Calibration retention", "Fail-safe defaults"],
    )
    prs.save(output_path)
    print(f"Generated {output_path}")


def _generate_eml() -> None:
    mail_dir = INPUT_DIR / "mail"
    mail_dir.mkdir(parents=True, exist_ok=True)
    emails = [
        {
            "filename": "mail1.eml",
            "subject": "Brake feature confirmation",
            "body": "Brake safety checks pass and error handling is stable.",
        },
        {
            "filename": "mail2.eml",
            "subject": "CAN timing diagnostics",
            "body": "CAN diagnostics confirm timing limits and calibration updates.",
        },
        {
            "filename": "mail3.eml",
            "subject": "Error handling summary",
            "body": "Error handling includes retries and logging in NVM.",
        },
    ]
    for idx, mail in enumerate(emails, start=1):
        content = (
            f"From: demo{idx}@example.com\n"
            "To: team@example.com\n"
            "Date: Fri, 01 Mar 2024 10:00:00 +0000\n"
            f"Subject: {mail['subject']}\n"
            f"Message-ID: <demo-{idx}@example.com>\n"
            "Content-Type: text/plain; charset=\"utf-8\"\n"
            "\n"
            f"{mail['body']}\n"
        )
        (mail_dir / mail["filename"]).write_text(content, encoding="utf-8")
    print(f"Generated {len(emails)} EML files in {mail_dir}")


if __name__ == "__main__":
    sys.exit(main())
