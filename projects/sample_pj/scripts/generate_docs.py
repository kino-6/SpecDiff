#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from openpyxl import Workbook
from pptx import Presentation


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = PROJECT_ROOT / "docs"


def build_pdf(path: Path, title: str, lines: list[str]) -> None:
    content_lines = [f"({title}) Tj"]
    for line in lines:
        safe = line.replace("\\", r"\\").replace("(", r"\\(").replace(")", r"\\)")
        content_lines.append(f"0 -18 Td ({safe}) Tj")
    content = "BT /F1 12 Tf 72 760 Td {} ET".format(" ".join(content_lines))
    content_bytes = content.encode("utf-8")

    objects = []

    def add_obj(data: bytes) -> int:
        objects.append(data)
        return len(objects)

    add_obj(b"<< /Type /Catalog /Pages 2 0 R >>")
    add_obj(b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>")
    add_obj(
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R "
        b"/Resources << /Font << /F1 5 0 R >> >> >>"
    )
    add_obj(b"<< /Length %d >>\nstream\n%s\nendstream" % (len(content_bytes), content_bytes))
    add_obj(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    xref_positions = []
    output = [b"%PDF-1.4\n"]
    for idx, obj in enumerate(objects, start=1):
        xref_positions.append(sum(len(chunk) for chunk in output))
        output.append(f"{idx} 0 obj\n".encode("ascii"))
        output.append(obj + b"\nendobj\n")

    xref_start = sum(len(chunk) for chunk in output)
    output.append(b"xref\n")
    output.append(f"0 {len(objects)+1}\n".encode("ascii"))
    output.append(b"0000000000 65535 f \n")
    for pos in xref_positions:
        output.append(f"{pos:010d} 00000 n \n".encode("ascii"))
    output.append(b"trailer\n")
    output.append(f"<< /Size {len(objects)+1} /Root 1 0 R >>\n".encode("ascii"))
    output.append(b"startxref\n")
    output.append(f"{xref_start}\n".encode("ascii"))
    output.append(b"%%EOF\n")

    path.write_bytes(b"".join(output))


def build_pdfs() -> None:
    spec_dir = DOCS_DIR / "spec"
    spec_dir.mkdir(parents=True, exist_ok=True)

    system_lines = [
        "The brake controller SHALL initialize comms within 100ms of init.",
        "The CAN interface MUST transmit brake status every 50ms timing window.",
        "Diagnostics SHALL record error_handling events for safety audits.",
        "Calibration values MUST be loaded from NVM before brake enable.",
        "The system SHALL enter safe mode on CAN bus errors.",
        "Brake pressure commands SHALL be limited to BRAKE_MAX_PRESSURE.",
        "The comms watchdog SHALL raise diagnostics on missed timing.",
        "The init sequence SHALL verify safety interlock before apply.",
    ]

    safety_lines = [
        "Safety diagnostics SHALL flag over-temperature brake faults.",
        "Error_handling MUST log calibration mismatches in NVM.",
        "The brake controller SHALL debounce timing jitter below 5ms.",
        "CAN comms MUST reject frames with invalid length.",
        "The system SHALL allow manual reset after safety latch.",
        "Diagnostics SHALL report active error codes on demand.",
        "Calibration SHALL be stored with NVM signature verification.",
        "Init SHALL confirm safety interlock on power-up.",
    ]

    build_pdf(spec_dir / "01_system_spec.pdf", "System Spec", system_lines)
    build_pdf(spec_dir / "02_safety_spec.pdf", "Safety Spec", safety_lines)


def build_xlsx() -> None:
    qa_dir = DOCS_DIR / "qa"
    qa_dir.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = "Q&A"
    ws.append(["Question", "Answer", "Status", "Owner"])
    rows = [
        ("How does brake timing sync with CAN comms?", "Timing aligns to 50ms cycle.", "Approved", "A. Lopez"),
        ("What diagnostics run at init?", "Self-test and safety interlock checks.", "Draft", "B. Chen"),
        ("How is calibration stored in NVM?", "Pressure offsets saved with signature.", "Approved", "C. Patel"),
        ("What is the error_handling response to CAN bus loss?", "Enter safe mode and log error.", "Approved", "D. Silva"),
        ("When is brake release commanded?", "On comms timeout or manual request.", "Draft", "E. Stein"),
        ("How is diagnostics data reported?", "CAN diagnostics frames every 100ms.", "Approved", "F. Wu"),
        ("What safety limits apply to brake pressure?", "Clamp to BRAKE_MAX_PRESSURE.", "Approved", "G. Rao"),
        ("How does init validate NVM?", "Checks signature before calibration.", "Draft", "H. Kim"),
        ("What retry policy is used for comms?", "Three retries with timing backoff.", "Approved", "I. Khan"),
        ("How are calibration overrides handled?", "Diagnostics must approve change.", "Draft", "J. Novak"),
        ("What CAN IDs carry brake status?", "0x120 for status, 0x121 for diag.", "Approved", "K. Ortiz"),
        ("How is safety interlock cleared?", "Manual reset after diagnostics pass.", "Approved", "L. Singh"),
    ]
    for row in rows:
        ws.append(row)

    wb.save(qa_dir / "qa_trace.xlsx")


def build_pptx() -> None:
    slides_dir = DOCS_DIR / "slides"
    slides_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    slides = [
        (
            "Design Review - Brake Comms",
            [
                "CAN timing budget: 50ms cycle",
                "Brake status comms includes safety flag",
                "Diagnostics frames report error_handling",
            ],
            "Notes: emphasize comms retry policy.",
        ),
        (
            "Safety & Diagnostics",
            [
                "Safety interlock before brake apply",
                "Diagnostics monitor temperature and CAN errors",
                "Calibration stored in NVM",
            ],
            "Notes: mention NVM signature check.",
        ),
        (
            "Calibration Workflow",
            [
                "Calibrate brake pressure offset",
                "Store calibration in NVM",
                "Init loads calibration before enable",
            ],
            None,
        ),
        (
            "Initialization Sequence",
            [
                "Init self-test",
                "Comms bus check",
                "Timing sync and watchdog",
            ],
            None,
        ),
        (
            "Error Handling",
            [
                "Retry CAN comms on timeout",
                "Enter safe mode on persistent errors",
                "Log diagnostics for audit",
            ],
            "Notes: highlight safety implications.",
        ),
    ]

    for title, bullets, notes in slides:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for bullet in bullets:
            p = body.add_paragraph()
            p.text = bullet
            p.level = 0
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(slides_dir / "design_review.pptx")


def main() -> None:
    build_pdfs()
    build_xlsx()
    build_pptx()


if __name__ == "__main__":
    main()
